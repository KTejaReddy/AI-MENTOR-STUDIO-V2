import os
import json
import uuid
import shutil
import hashlib
import asyncio
from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks, Depends, Request
from fastapi.responses import StreamingResponse

# Lazy loaded document libraries
_doc_libs_cache = None

def get_document_libraries():
    global _doc_libs_cache
    if _doc_libs_cache is None:
        import fitz  # PyMuPDF
        import docx
        import pptx
        _doc_libs_cache = (fitz, docx, pptx)
    return _doc_libs_cache

from app.core.rate_limit import limiter
from app.core.dependencies import get_current_user
from app.models.user import User

from app.ai.gateway import gateway
from app.ai.key_manager import key_manager
from app.ai.groq_provider import GroqProvider
from app.schemas.document import DocumentUploadResponse

import time
import logging
import traceback
logger = logging.getLogger(__name__)

router = APIRouter()

DOCS_DIR = os.path.join(os.getcwd(), "data", "documents")
os.makedirs(DOCS_DIR, exist_ok=True)

CACHE_DIR = os.path.join(DOCS_DIR, "cache")
os.makedirs(CACHE_DIR, exist_ok=True)

def extract_text(file_path: str, filename: str) -> str:
    logger.info(f"Text extraction start for: {filename}")
    start_time = time.time()
    ext = filename.split('.')[-1].lower()
    text = ""
    try:
        fitz, docx, pptx = get_document_libraries()
        if ext == "pdf":
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n\n"
        elif ext in ("docx", "doc"):
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in ("pptx", "ppt"):
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse document: {str(e)}")
    
    elapsed = time.time() - start_time
    logger.info(f"Text extraction complete in {elapsed:.2f}s")
    return text

async def _run_document_analysis(doc_id: str, filename: str, file_hash: str, user_id: str):
    logger.info(f"Starting explanation generation for document_id: {doc_id}")
    start_time = time.time()
    
    user_dir = os.path.join(DOCS_DIR, user_id)
    os.makedirs(user_dir, exist_ok=True)
    file_path = os.path.join(user_dir, f"{doc_id}_file")
    json_path = os.path.join(user_dir, f"{doc_id}.json")
    
    try:
        text = extract_text(file_path, filename)
        if not text.strip():
            raise ValueError("Document contains no readable text.")
            
        logger.info(f"Explanation generation start for {doc_id}")
        
        prompt = f"""You are an excellent university teacher. Read the uploaded document completely.
Explain the content in simple language. Assume the reader is learning the topic for the first time.
Never invent information. Explain only what exists inside the uploaded document. Preserve technical accuracy.
Avoid unnecessary repetition. Use proper Markdown.

OUTPUT FORMAT:
Generate the explanation using Markdown.
Structure:
# Document Overview
Briefly explain what the document is about.
---
# Main Concepts
Explain every important concept using simple language.
---
# Important Definitions
Explain important terms.
---
# Key Formulas
Only if formulas exist. Explain what each formula means.
---
# Examples
Provide simple examples based only on the document. Do not invent new topics.
---
# Important Points
List the most important ideas.
---
# Summary
Provide a concise revision summary.

Document Text (truncated if necessary):
{text[:40000]}"""
        provider = GroqProvider(key_manager)
        response = await provider.complete(
            type("CompletionRequest", (object,), {
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": prompt}],
                "to_dict": lambda s: {"model": s.model, "messages": s.messages}
            })()
        )
        
        explanation = response.content
        
        doc_data = {
            "id": doc_id,
            "filename": filename,
            "status": "ready",
            "file_hash": file_hash,
            "explanation": explanation
        }
        
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(doc_data, f)
            
        # Cache globally
        cache_json_path = os.path.join(CACHE_DIR, f"{file_hash}.json")
        with open(cache_json_path, "w", encoding="utf-8") as f:
            json.dump(doc_data, f)
        
        logger.info(f"Total analysis completed in {time.time() - start_time:.2f}s for {doc_id}")
        return doc_data
        
    except Exception as e:
        logger.error(f"Error in document explanation for {doc_id}: {e}\n{traceback.format_exc()}")
        try:
            doc_data = {
                "id": doc_id,
                "filename": filename,
                "status": "failed",
                "file_hash": file_hash,
                "error": str(e),
                "explanation": ""
            }
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(doc_data, f)
        except Exception:
            pass
        raise

@router.get("/list")
async def list_documents(current_user: User = Depends(get_current_user)):
    user_dir = os.path.join(DOCS_DIR, current_user.id)
    if not os.path.exists(user_dir):
        return {"documents": []}
        
    docs = []
    for f in os.listdir(user_dir):
        if f.endswith(".json") and f != "cache":
            try:
                with open(os.path.join(user_dir, f), "r", encoding="utf-8") as file:
                    data = json.load(file)
                    docs.append({
                        "document_id": data.get("id"),
                        "filename": data.get("filename"),
                        "title": data.get("filename"),
                        "status": data.get("status", "ready"),
                    })
            except:
                pass
    return {"documents": docs}

@router.get("/{document_id}")
async def get_document(document_id: str, current_user: User = Depends(get_current_user)):
    user_dir = os.path.join(DOCS_DIR, current_user.id)
    file_path = os.path.join(user_dir, f"{document_id}.json")
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Document not found")
    with open(file_path, "r", encoding="utf-8") as f:
        return json.load(f)

@router.delete("/{document_id}")
async def delete_document(document_id: str, current_user: User = Depends(get_current_user)):
    user_dir = os.path.join(DOCS_DIR, current_user.id)
    for ext in [".json", "_file"]:
        path = os.path.join(user_dir, f"{document_id}{ext}")
        if os.path.exists(path):
            try:
                os.remove(path)
            except Exception:
                pass
    return {"status": "deleted"}

@router.post("/upload", response_model=DocumentUploadResponse)
async def upload_document(background_tasks: BackgroundTasks, file: UploadFile = File(...), current_user: User = Depends(get_current_user)):
    logger.info("Upload start")
    
    ALLOWED_MIME_TYPES = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "text/plain"
    }
    
    if file.content_type not in ALLOWED_MIME_TYPES:
        raise HTTPException(status_code=400, detail="Invalid file type. Only PDF, DOCX, PPTX, and TXT are allowed.")
        
    doc_id = str(uuid.uuid4())
    user_dir = os.path.join(DOCS_DIR, current_user.id)
    os.makedirs(user_dir, exist_ok=True)
    temp_file_path = os.path.join(user_dir, f"{doc_id}_file")
    
    # Read bytes to compute hash
    file_bytes = await file.read()
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    
    cache_json_path = os.path.join(CACHE_DIR, f"{file_hash}.json")
    
    if os.path.exists(cache_json_path):
        logger.info("Cache hit: reusing previously parsed document analysis")
        with open(cache_json_path, "r", encoding="utf-8") as f:
            cached_data = json.load(f)
            
        cached_data["id"] = doc_id
        with open(os.path.join(user_dir, f"{doc_id}.json"), "w", encoding="utf-8") as f:
            json.dump(cached_data, f)
            
        # Save dummy empty file for delete/list requirements
        with open(temp_file_path, "wb") as f:
            f.write(b"")
            
        return DocumentUploadResponse(
            document_id=doc_id,
            filename=file.filename,
            pages=1,
            status="uploaded"
        )
        
    # Write actual file contents
    with open(temp_file_path, "wb") as f:
        f.write(file_bytes)
        
    # Initialize placeholder json
    doc_data = {
        "id": doc_id,
        "filename": file.filename,
        "status": "processing",
        "file_hash": file_hash,
        "explanation": ""
    }
    with open(os.path.join(user_dir, f"{doc_id}.json"), "w", encoding="utf-8") as f:
        json.dump(doc_data, f)
        
    # Trigger background parsing task
    background_tasks.add_task(_run_document_analysis, doc_id, file.filename, file_hash, current_user.id)
    
    return DocumentUploadResponse(
        document_id=doc_id,
        filename=file.filename,
        pages=1,
        status="uploaded"
    )
